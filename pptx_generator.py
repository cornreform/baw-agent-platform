#!/usr/bin/env python3
"""
PPTX Generator - Create PowerPoint presentations using python-pptx
Supports text slides and embedded chart images
"""

import json
import os
import sys
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# ANSI colors for terminal output
class Colors:
    GREEN = '\033[92m'
    YELLOW = '\033[93m'
    BLUE = '\033[94m'
    RESET = '\033[0m'

def load_slide_content(json_file="slide_content.json"):
    """Load slide content from JSON file"""
    if not os.path.exists(json_file):
        print(f"No slide content JSON file found: {json_file}")
        return None
    
    with open(json_file, 'r') as f:
        return json.load(f)

def create_presentation(content, output_file="presentation.pptx"):
    """Create PPTX presentation from content"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 widescreen
    prs.slide_height = Inches(7.5)
    
    # Add title slide first
    title_slide_layout = prs.slide_layouts[0]
    title_slide = prs.slides.add_slide(title_slide_layout)
    title_slide.shapes.title.text = content.get('title', 'Presentation')
    title_slide.placeholders[1].text = content.get('subtitle', 'Generated with python-pptx')
    
    slide_count = 0
    
    for slide_data in content.get('slides', []):
        slide_type = slide_data.get('type', 'text')
        
        if slide_type == 'text':
            # Create title + content slide
            slide_layout = prs.slide_layouts[1]  # Title and Content
            slide = prs.slides.add_slide(slide_layout)
            slide_count += 1
            
            # Title
            title = slide.shapes.title
            title.text = slide_data.get('title', 'Untitled')
            
            # Content
            content_box = slide.placeholders[1]
            tf = content_box.text_frame
            tf.clear()
            
            for line in slide_data.get('content', []):
                p = tf.add_paragraph()
                p.text = line
                p.font.size = Pt(18)
                p.space_after = Pt(8)
                
        elif slide_type == 'chart':
            # Create slide with chart image
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            slide_count += 1
            
            # Title
            title = slide.shapes.title
            title.text = slide_data.get('title', 'Chart')
            
            # Add chart image
            chart_file = slide_data.get('chart_file', 'chart_sample.png')
            if os.path.exists(chart_file):
                # Add image to slide
                left = Inches(1)
                top = Inches(1.5)
                width = Inches(6)
                slide.shapes.add_picture(chart_file, left, top, width=width)
                
                # Add data table on the right
                chart_info = slide_data.get('chart_data', {})
                labels = chart_info.get('labels', [])
                values = chart_info.get('values', [])
                
                # Create textbox for data
                left = Inches(7.5)
                top = Inches(1.5)
                width = Inches(5)
                height = Inches(4)
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.word_wrap = True
                
                p = tf.add_paragraph()
                p.text = "Data:"
                p.font.size = Pt(14)
                p.font.bold = True
                
                for label, value in zip(labels, values):
                    p = tf.add_paragraph()
                    p.text = f"{label}: {value}"
                    p.font.size = Pt(12)
                    p.space_after = Pt(4)
        
        elif slide_type == 'text_chart':
            # Create slide with both text and chart
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            slide_count += 1
            
            # Title
            title = slide.shapes.title
            title.text = slide_data.get('title', 'Analysis')
            
            # Add text on left
            left = Inches(0.5)
            top = Inches(1.5)
            width = Inches(4)
            height = Inches(5)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True
            
            for line in slide_data.get('content', []):
                p = tf.add_paragraph()
                p.text = line
                p.font.size = Pt(14)
                p.space_after = Pt(6)
            
            # Add chart image on right
            chart_file = slide_data.get('chart_file', 'chart_sample.png')
            if os.path.exists(chart_file):
                left = Inches(5)
                top = Inches(1.5)
                width = Inches(7.5)
                slide.shapes.add_picture(chart_file, left, top, width=width)
    
    # Save presentation
    prs.save(output_file)
    return slide_count

def main():
    print(f"{Colors.BLUE}=== PPTX Generator ==={Colors.RESET}")
    
    # Load content from JSON file
    content = load_slide_content()
    
    if content is None:
        # Create sample content
        print(f"{Colors.YELLOW}Creating sample presentation...{Colors.RESET}")
        content = {
            "title": "Sample Presentation",
            "subtitle": "Generated with python-pptx",
            "slides": [
                {
                    "type": "text",
                    "title": "Overview",
                    "content": [
                        "This is a sample presentation created with python-pptx.",
                        "It demonstrates text, charts, and combined layouts.",
                        "The library provides powerful PPTX generation capabilities."
                    ]
                },
                {
                    "type": "chart",
                    "title": "Sales Data",
                    "chart_file": "chart_bar.png",
                    "chart_data": {
                        "labels": ["Q1", "Q2", "Q3", "Q4"],
                        "values": [100, 150, 120, 180]
                    }
                },
                {
                    "type": "text_chart",
                    "title": "Performance Analysis",
                    "content": [
                        "Q4 showed the highest growth.",
                        "Overall trend is positive.",
                        "Target exceeded by 20%."
                    ],
                    "chart_file": "chart_line.png"
                }
            ]
        }
    
    # Create presentation
    output_file = "output.pptx"
    count = create_presentation(content, output_file)
    
    print(f"{Colors.GREEN}✓ Created {output_file} with {count} slides{Colors.RESET}")
    print(f"File size: {os.path.getsize(output_file):,} bytes")

if __name__ == "__main__":
    main()