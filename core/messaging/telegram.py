"""
BAW — Telegram Bot Connector

Long-polling Telegram Bot via httpx (no extra dependencies).
Fully featured: commands, replies, error handling, reconnection.
"""
from __future__ import annotations
import json
import logging
import os
import threading
import time
import httpx
from pathlib import Path
from typing import Optional

from . import BaseConnector, Message, register

logger = logging.getLogger("baw.telegram")

POLL_TIMEOUT = 30  # Long-poll timeout (seconds)
POLL_RETRY_DELAY = 5
MAX_MESSAGE_LENGTH = 4000
API_BASE = "https://api.telegram.org/bot{token}"


@register("telegram", "Telegram Bot — long-polling via httpx", "telegram")
class TelegramConnector(BaseConnector):
    """Telegram Bot connector using long-polling (getUpdates).

    Config:
      telegram:
        token: "***"          # Bot token from @BotFather
        allowed_users: []     # Optional: list of user IDs to allow
    """

    def __init__(self, config: dict, on_message):
        super().__init__(config.get("telegram", {}), on_message)
        self._token = os.environ.get("TELEGRAM_BOT_TOKEN") or self.config.get("token", "")
        self._allowed = [str(u) for u in self.config.get("allowed_users", [])]
        self._offset = 0
        self._client: httpx.Client | None = None
        self._api_base = API_BASE.format(token=self._token) if self._token else ""
        self._debounce_until = 0.0
        self._restart_chat_id: str | None = None
        self._tts_enabled = self.config.get("tts_enabled", False)
        self._tts_voice = self.config.get("tts_voice", "male-tone-1")
        self._selector_msg_id: int | None = None
        self._selector_role: dict[str, str] = {}  # {chat_id: role} for 3-layer model selector
        # ── Media group buffering ──
        self._media_group_buffers: dict[str, dict[str, list]] = {}  # {chat_id: {group_id: [msgs]}}
        self._media_group_timers: dict[str, threading.Timer] = {}
        self._media_group_wait = 2.5  # seconds to wait for all media in a group

    def connect(self) -> bool:
        """Test connection by fetching bot info."""
        if not self._token:
            logger.error("[Telegram] No token configured")
            return False
        try:
            self._client = httpx.Client(timeout=10)
            r = self._client.get(f"{self._api_base}/getMe")
            if r.status_code == 200:
                info = r.json()
                if info.get("ok"):
                    bot_name = info["result"].get("first_name", "BAW Bot")
                    logger.info(f"[Telegram] Connected as @{info['result'].get('username', '?')}")

                    # Register slash command menu
                    self._register_commands()

                    # ── Key Vault: backup .env keys to secure storage ──
                    try:
                        from ..key_vault import backup as _kv_backup, restore as _kv_restore
                        _kv_restore()  # auto-restore if .env is missing keys
                        _kv_backup()   # sync vault with current .env
                    except Exception as _kve:
                        logger.warning(f"[KeyVault] init failed: {_kve}")

                    # ── Provider health ping (startup check) + dead-provider notification ──
                    try:
                        from ..llm import ping_provider_health
                        from ..config import load_config as _load_full_cfg
                        _full_cfg = _load_full_cfg(reload=True)
                        _health = ping_provider_health(_full_cfg)
                        _dead = {k: v for k, v in _health.items() if v not in ("healthy", "key_set", "no_key", "no_key_config", "auth_error")}
                        if _dead:
                                logger.warning(f"[Telegram] Dead providers at startup: {list(_dead.keys())}")
                                
                                # Build user notification with fix suggestions
                                _dead_list = "\n".join(f"  • {k}: {v}" for k, v in _dead.items())
                                _fix_lines = []
                                
                                # Check if any dead provider is referenced in config
                                _def_model = _full_cfg.get("model", {}).get("default", "")
                                _fallback = _full_cfg.get("model", {}).get("fallback", "")
                                _caps = _full_cfg.get("capabilities", {})
                                
                                for _pname in _dead:
                                    # Check which models on this provider are in use
                                    _models = [m.get("id","?") for m in _full_cfg.get("providers",{}).get(_pname,{}).get("models",[])]
                                    
                                    if _def_model in _models:
                                        _fix_lines.append(f"  /set model.default deepseek-v4-flash  ← default 指向死 provider「{_pname}」")
                                    if _fallback in _models:
                                        _fix_lines.append(f"  /set model.fallback deepseek-v4-pro  ← fallback 指向死 provider「{_pname}」")
                                    for _cap, _cc in _caps.items():
                                        if isinstance(_cc, dict) and _cc.get("model") in _models:
                                            _fix_lines.append(f"  /set capabilities.{_cap}.model deepseek-v4-flash  ← {_cap} 指向死 provider「{_pname}」")
                                
                                _fix_section = "\n".join(_fix_lines) if _fix_lines else "  冇 config reference — 毋須改動"
                                
                                _admin_id = os.environ.get("BAW_ADMIN_CHAT_ID", "")
                                _notify_text = (
                                    f"<b>⚠️  Provider Health Alert</b>\n\n"
                                    f"以下 provider 無法連線：\n{_dead_list}\n\n"
                                    f"<b>建議修改：</b>\n{_fix_section}\n\n"
                                    f"系統已自動 blacklist 死 provider，"
                                    f"BAW 會自動 fallback 到可用 provider。"
                                )
                                if _admin_id:
                                    try:
                                        self._client.post(
                                            f"{self._api_base}/sendMessage",
                                            json={"chat_id": _admin_id, "text": _notify_text, "parse_mode": "HTML"},
                                            timeout=10,
                                        )
                                    except Exception:
                                        pass
                    except Exception as _he:
                        logger.warning(f"[Telegram] Provider health ping failed: {_he}")

                    # ── Back-online notification after restart ──
                    self._notify_restart()

                    return True
            logger.error(f"[Telegram] getMe failed: {r.status_code} {r.text[:200]}")
            # ── User-friendly fatal error notification ──
            _err_msg = ""
            if r.status_code == 401:
                _err_msg = "<b>Telegram Bot 連接失敗</b>\n\n"
                _err_msg += "原因: Bot Token 無效或已被撤銷\n"
                _err_msg += "解決: 請去 @BotFather 重新生成 Token, 然後更新 ~/.baw/.env"
            elif r.status_code == 404:
                _err_msg = "<b>Telegram Bot 連接失敗</b>\n\n"
                _err_msg += "原因: API 端點不存在 (404)\n"
                _err_msg += "解決: 檢查網絡連線或 Telegram API 狀態"
            elif r.status_code == 429:
                _err_msg = "<b>Telegram 限流</b>\n\n"
                _err_msg += "原因: 發送太多請求被限流\n"
                _err_msg += "解決: 等幾分鐘後會自動恢復"
            else:
                _err_msg = f"<b>Telegram Bot 連接失敗</b>\n\n"
                _err_msg += f"原因: HTTP {r.status_code}\n"
                _err_msg += "解決: 檢查 .env 中的 TELEGRAM_BOT_TOKEN 是否正確"
            # Try to notify admin if we know the chat_id
            _admin_id = os.environ.get("BAW_ADMIN_CHAT_ID", "")
            if _admin_id:
                try:
                    self._client.post(
                        f"{self._api_base}/sendMessage",
                        json={"chat_id": _admin_id, "text": _err_msg, "parse_mode": "HTML"},
                        timeout=10,
                    )
                except Exception:
                    pass
            # Also write to a visible file for admin
            _fatal_path = Path.home() / ".baw" / "FATAL_ERROR.txt"
            _fatal_path.write_text(_err_msg, encoding="utf-8")
            return False
        except Exception as e:
            logger.error(f"[Telegram] Connection error: {e}")
            return False

    def _notify_restart(self):
        """Send 'Back Online' notification if this was a restart."""
        import json as _json
        from pathlib import Path
        pending_file = Path.home() / ".baw" / ".restart_pending"
        if not pending_file.exists():
            return
        try:
            data = _json.loads(pending_file.read_text())
            chat_id = data.get("chat_id", "")
            if chat_id:
                self._client.post(
                    f"{self._api_base}/sendMessage",
                    json={"chat_id": chat_id, "text": "<b>BAW Back Online</b>", "parse_mode": "HTML"},
                    timeout=10,
                )
                logger.info(f"[Telegram] Restart notification sent to {chat_id}")
            pending_file.unlink(missing_ok=True)
        except Exception as e:
            logger.warning(f"[Telegram] Restart notification failed: {e}")

    def _register_commands(self):
        """Register bot command menu via setMyCommands."""
        commands = [
            {"command": "start",   "description": "Welcome message"},
            {"command": "help",    "description": "Show all commands"},
            {"command": "status",  "description": "BAW system status + sessions"},
            {"command": "btw",     "description": "Quick answer (no court)"},
            {"command": "model",   "description": "Switch model or show model selector"},
            {"command": "mode",    "description": "Switch mode: quick / hybrid / tight"},
            {"command": "tone",    "description": "Switch tone: casual / business / teaching"},
            {"command": "set",     "description": "Persist config: /set key value"},
            {"command": "court",   "description": "Show last Angel/Devil verdict"},
            {"command": "fresh",   "description": "Raw model — no soul, no memories"},
            {"command": "memory",  "description": "Save a memory entry"},
            {"command": "search",  "description": "Search stored memories"},
            {"command": "board",   "description": "Generate HTML dashboard"},
            {"command": "task",    "description": "Session: new/list/resume/save/forget/info"},
            {"command": "list",    "description": "List saved sessions (alias for /task list)"},
            {"command": "new",     "description": "Save current + start fresh session"},
            {"command": "reset",   "description": "Hard reset — clear session without saving"},
            {"command": "resume",  "description": "Resume a saved session"},
            {"command": "summarize", "description": "LLM summary of current session"},
            {"command": "compact",   "description": "Manual compress: summarize early convo + keep last 4 msgs"},
            {"command": "pickup",  "description": "Resume last interrupted session"},
            {"command": "reload",  "description": "Hot-reload tools & config"},
            {"command": "evolve",  "description": "Self-evolution stats"},
            {"command": "tts",     "description": "Toggle TTS: on / off / status"},
            {"command": "capability", "description": "Manage capabilities"},
            {"command": "update",  "description": "Git pull + changelog + restart"},
            {"command": "stop",    "description": "Cancel running request"},
            {"command": "restart", "description": "Restart BAW engine"},
            {"command": "doctor",  "description": "10-point system health check"},
            {"command": "backup",  "description": "Create/list/restore backups"},
            {"command": "monitor", "description": "24h error rate + weekly report"},
            {"command": "queue",   "description": "Show pending message queue"},
        ]
        try:
            self._client.post(
                f"{self._api_base}/setMyCommands",
                json={"commands": commands},
                timeout=10,
            )
        except Exception as e:
            logger.warning(f"[Telegram] Failed to register commands: {e}")

    def disconnect(self):
        if self._client:
            self._client.close()
            self._client = None

    def send(self, chat_id: str, text: str, edit_msg_id: str = "") -> str:
        """Send or edit a message. Returns message_id if successful, empty string if failed.
        If edit_msg_id is provided, edits that message instead of sending new one.
        
        Long messages (>2000 chars) are auto-truncated to the first section + summary.
        Set display.output_mode: full in config to disable truncation."""
        if not self._client or not self._token:
            return ""
        # ── Intercept model selector ——
        if text.startswith("[MODEL_ROLE_SELECT]\n"):
            self._send_role_selector(chat_id, text)
            return ""
        if text.startswith("[MODEL_SELECT]\n"):
            self._send_model_selector_text(chat_id, text)
            return ""
        try:
            # ── Extract MEDIA: tags ──
            import re as _re
            media_files = _re.findall(r'^MEDIA:(.+)$', text, _re.MULTILINE)
            if media_files:
                text = _re.sub(r'^MEDIA:.+$\n?', '', text, flags=_re.MULTILINE).strip()
                # Can't edit media — send new message
                edit_msg_id = ""

            # ── Edit existing message ──
            if edit_msg_id:
                if text:
                    return self._edit_text(chat_id, edit_msg_id, text)
                return edit_msg_id

            # ── Send new text ──
            msg_id = ""
            if text:
                # Telegram max 4096 chars — split into multiple messages if needed
                if len(text) > MAX_MESSAGE_LENGTH:
                    parts = []
                    remaining = text
                    while remaining:
                        parts.append(remaining[:MAX_MESSAGE_LENGTH])
                        remaining = remaining[MAX_MESSAGE_LENGTH:]
                    for part in parts:
                        msg_id = self._send_text(chat_id, part)
                else:
                    msg_id = self._send_text(chat_id, text)

            # ── Send media files ──
            for fpath in media_files:
                fpath = fpath.strip()
                # ── Container→host file resolution ──
                fpath = self._resolve_media_path(fpath)
                if not fpath:
                    continue
                self._send_media(chat_id, fpath)

            return msg_id or ""
        except Exception as e:
            logger.error(f"[Telegram] send error: {e}")
            return ""

    def _resolve_media_path(self, fpath: str) -> str:
        """Resolve a MEDIA file path. Verifies the file actually exists
        and is readable. If not, logs and returns empty string so the
        caller can skip the attachment (instead of sending 'File not found'
        error to the user)."""
        from pathlib import Path
        fpath_obj = Path(fpath)
        if fpath_obj.exists() and fpath_obj.is_file():
            return str(fpath_obj.resolve())

        # Last-ditch: try /tmp/<basename> and /home/baw/.baw/media/tts/<basename>
        # in case the path was hallucinated by sub-agent planning
        basename = fpath_obj.name
        for _candidate in [
            Path(f"/tmp/{basename}"),
            Path(f"/home/baw/.baw/media/tts/{basename}"),
            Path(f"/home/baw/.baw/{basename}"),
        ]:
            if _candidate.exists() and _candidate.is_file():
                logger.info(f"[Telegram] MEDIA path fallback: {fpath} → {_candidate}")
                return str(_candidate.resolve())

        # File genuinely not found — log clearly so the diagnostic shows up
        logger.warning(f"[Telegram] MEDIA file not found anywhere: {fpath}")
        return ""

    def _send_text(self, chat_id: str, text: str) -> str:
        """Send a plain text message. Returns message_id string or empty on failure."""
        r = self._client.post(
            f"{self._api_base}/sendMessage",
            json={"chat_id": chat_id, "text": text, "parse_mode": "HTML"},
            timeout=10,
        )
        if r.status_code != 200:
            if "can't parse entities" in r.text:
                # Strip all HTML tags so raw <b> isn't shown literally
                import re as _re
                _clean = _re.sub(r'<[^>]+>', '', text)
                r = self._client.post(
                    f"{self._api_base}/sendMessage",
                    json={"chat_id": chat_id, "text": _clean},
                    timeout=10,
                )
        if r.status_code == 200:
            data = r.json()
            if data.get("ok"):
                return str(data["result"]["message_id"])
        return ""

    def _edit_text(self, chat_id: str, message_id: str, text: str) -> str:
        """Edit an existing message. Returns message_id on success, empty on failure."""
        r = self._client.post(
            f"{self._api_base}/editMessageText",
            json={"chat_id": chat_id, "message_id": int(message_id), "text": text, "parse_mode": "HTML"},
            timeout=10,
        )
        if r.status_code != 200:
            if "can't parse entities" in r.text:
                r = self._client.post(
                    f"{self._api_base}/editMessageText",
                    json={"chat_id": chat_id, "message_id": int(message_id), "text": text},
                    timeout=10,
                )
        if r.status_code == 200:
            return message_id
        # Telegram returns 400 "message is not modified" when content is identical —
        # this is not an error, the message already shows the right content.
        if "message is not modified" in r.text:
            return message_id
        logger.warning(f"[Telegram] editMessageText failed: {r.status_code} {r.text[:200]}")
        return ""

    def send_typing(self, chat_id: str) -> bool:
        """Send typing indicator to show '...' in Telegram."""
        if not self._client or not self._token:
            return False
        try:
            r = self._client.post(
                f"{self._api_base}/sendChatAction",
                json={"chat_id": chat_id, "action": "typing"},
                timeout=5,
            )
            return r.status_code == 200
        except Exception:
            return False

    _MEDIA_EXT_MAP = {
        ".png": "sendPhoto", ".jpg": "sendPhoto", ".jpeg": "sendPhoto",
        ".webp": "sendPhoto", ".gif": "sendPhoto",
        ".mp4": "sendVideo", ".mov": "sendVideo",
        ".mp3": "sendAudio", ".wav": "sendAudio",
        ".ogg": "sendVoice",
    }

    def _send_media(self, chat_id: str, file_path: str):
        """Send a media file using the appropriate Telegram API endpoint."""
        from pathlib import Path

        fpath = Path(file_path).expanduser()
        if not fpath.exists():
            logger.warning(f"[Telegram] MEDIA file not found: {fpath}")
            self._send_text(chat_id, f"File not found: {fpath}")
            return

        ext = fpath.suffix.lower()
        endpoint = self._MEDIA_EXT_MAP.get(ext, "sendDocument")

        field_map = {
            "sendPhoto": "photo", "sendDocument": "document",
            "sendVideo": "video", "sendAudio": "audio", "sendVoice": "voice",
        }
        field = field_map.get(endpoint, "document")

        try:
            with open(fpath, "rb") as f:
                files = {field: (fpath.name, f)}
                r = self._client.post(
                    f"{self._api_base}/{endpoint}",
                    data={"chat_id": chat_id},
                    files=files,
                    timeout=60,
                )
            if r.status_code != 200:
                logger.warning(f"[Telegram] Media send failed ({endpoint}): {r.text[:200]}")
                if endpoint != "sendDocument":
                    self._send_media_as_document(chat_id, fpath)
        except Exception as e:
            logger.error(f"[Telegram] Media send error ({fpath}): {e}")

    def _send_media_as_document(self, chat_id: str, fpath):
        """Fallback: send any file as a document."""
        from pathlib import Path
        fpath = Path(fpath)
        if not fpath.exists():
            return
        try:
            with open(fpath, "rb") as f:
                self._client.post(
                    f"{self._api_base}/sendDocument",
                    data={"chat_id": chat_id},
                    files={"document": (fpath.name, f)},
                    timeout=60,
                )
        except Exception as e:
            logger.error(f"[Telegram] Document fallback error ({fpath}): {e}")

    # ── File processing ─────────────────────────────────────────

    def _download_file(self, file_id: str, file_name: str = "file") -> str:
        """Download a file from Telegram Bot API. Returns local path."""
        import os
        from pathlib import Path

        # Get file path from Telegram
        r = self._client.post(
            f"{self._api_base}/getFile",
            json={"file_id": file_id},
            timeout=10,
        )
        if r.status_code != 200:
            err_text = r.text[:300]
            if "file is too big" in err_text.lower():
                raise RuntimeError(
                    f"📁 檔案太大，Telegram Bot API 限制 20MB。請用其他方式傳送（如壓縮、分割、或改用雲端連結）"
                )
            raise RuntimeError(f"getFile failed: {err_text}")
        data = r.json()
        if not data.get("ok"):
            err_desc = str(data.get("description", ""))
            if "file is too big" in err_desc.lower():
                raise RuntimeError(
                    f"📁 檔案太大，Telegram Bot API 限制 20MB。請用其他方式傳送（如壓縮、分割、或改用雲端連結）"
                )
            raise RuntimeError(f"getFile returned error: {data}")
        tg_path = data["result"]["file_path"]

        # Download file
        dl_url = f"https://api.telegram.org/file/bot{self._token}/{tg_path}"
        dl_r = self._client.get(dl_url, timeout=120)
        if dl_r.status_code != 200:
            raise RuntimeError(f"File download failed: {dl_r.status_code}")

        # Save to temp
        tmp_dir = Path.home() / ".baw" / "downloads"
        tmp_dir.mkdir(parents=True, exist_ok=True)
        local_path = tmp_dir / file_name
        # Avoid collisions
        if local_path.exists():
            base = local_path.stem
            ext = local_path.suffix
            i = 1
            while local_path.exists():
                local_path = tmp_dir / f"{base}_{i}{ext}"
                i += 1
        local_path.write_bytes(dl_r.content)
        logger.info(f"[Telegram] Downloaded: {tg_path} → {local_path} ({len(dl_r.content)} bytes)")
        return str(local_path)

    def _extract_file_content(self, file_path: str) -> str:
        """Extract text content from a file based on extension. Returns str."""
        from pathlib import Path

        fp = Path(file_path)
        ext = fp.suffix.lower()
        name = fp.name

        # ── Plain text ──
        if ext in (".txt", ".md", ".csv", ".json", ".xml", ".yaml", ".yml",
                    ".log", ".ini", ".cfg", ".conf", ".toml", ".py",
                    ".js", ".ts", ".html", ".css", ".sh", ".bash"):
            try:
                return fp.read_text(encoding="utf-8")
            except UnicodeDecodeError:
                return fp.read_text(encoding="latin-1")

        # ── PowerPoint ──
        elif ext in (".pptx", ".pptm"):
            try:
                from pptx import Presentation
                prs = Presentation(fp)
                parts = []
                for i, slide in enumerate(prs.slides, 1):
                    texts = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            texts.append(shape.text.strip())
                    if texts:
                        parts.append(f"--- Slide {i} ---\n" + "\n".join(texts))
                return "\n\n".join(parts) if parts else "[No text found in slides]"
            except Exception as e:
                return f"[Error extracting PowerPoint: {e}]"

        # ── Word ──
        elif ext in (".docx", ".docm"):
            try:
                from docx import Document
                doc = Document(fp)
                paras = [p.text for p in doc.paragraphs if p.text.strip()]
                return "\n\n".join(paras) if paras else "[No text found in document]"
            except Exception as e:
                return f"[Error extracting Word: {e}]"

        # ── PDF ──
        elif ext == ".pdf":
            try:
                # Try PyMuPDF first (fastest)
                import fitz
                doc = fitz.open(fp)
                pages = []
                for i, page in enumerate(doc, 1):
                    text = page.get_text().strip()
                    if text:
                        pages.append(f"--- Page {i} ---\n{text}")
                doc.close()
                result = "\n\n".join(pages)
                if result.strip():
                    return result
                # Fallback: pdftotext if installed
                import subprocess as sp
                txt_path = fp.with_suffix(".txt")
                sp.run(["pdftotext", str(fp), str(txt_path)], capture_output=True, timeout=30)
                if txt_path.exists():
                    return txt_path.read_text(encoding="utf-8")
                return "[No text extracted from PDF]"
            except Exception as e:
                return f"[Error extracting PDF: {e}]"

        # ── Excel ──
        elif ext in (".xlsx", ".xlsm"):
            try:
                import openpyxl
                wb = openpyxl.load_workbook(fp, read_only=True, data_only=True)
                parts = []
                for sheet_name in wb.sheetnames:
                    ws = wb[sheet_name]
                    rows = []
                    for row in ws.iter_rows(values_only=True):
                        vals = [str(v) if v is not None else "" for v in row]
                        line = " | ".join(vals)
                        if line.strip():
                            rows.append(line)
                    if rows:
                        parts.append(f"--- Sheet: {sheet_name} ---\n" + "\n".join(rows))
                wb.close()
                return "\n\n".join(parts) if parts else "[No data found in spreadsheet]"
            except Exception as e:
                return f"[Error extracting Excel: {e}]"

        # ── CSV (fallback for .csv not caught above) ──
        elif ext == ".csv":
            import csv, io
            try:
                text = fp.read_text(encoding="utf-8")
                reader = csv.reader(io.StringIO(text))
                lines = [" | ".join(row) for row in reader]
                return "\n".join(lines)
            except Exception as e:
                return f"[Error reading CSV: {e}]"

        # ── Images ──
        elif ext in (".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp", ".tiff"):
            try:
                import pytesseract
                from PIL import Image
                img = Image.open(fp)
                text = pytesseract.image_to_string(img, lang="chi_sim+eng")
                if text.strip():
                    return f"[OCR extracted text from {name}]\n\n{text}"
                return f"[Image {name}: no text detected via OCR]"
            except Exception as e:
                return f"[Error processing image {name}: {e}]"

        # ── Unknown ──
        else:
            return f"[Unsupported file type: {ext}. Cannot extract content.]"

    # ─────────────────────────────────────────────────────────────
    #  Media group handling + non-text message queueing
    # ─────────────────────────────────────────────────────────────

    def _handle_media_msg(self, chat_id: str, user_id, user_name: str, msg: dict, msg_type: str):
        """Single non-text message: acquire slot or queue (never drop)."""
        if self._acquire_slot():
            # Per-chat sequential: if this chat already has a task, queue it
            if self._is_chat_busy(chat_id):
                with self._active_lock:
                    self._active_count = max(0, self._active_count - 1)
                pos = self._enqueue_message(chat_id, user_id, user_name, "", msg, msg_type=msg_type)
                self.send(chat_id, f"⏳ Queued #{pos} — another task is running in this chat")
                return
            self._mark_chat_busy(chat_id)
            self._cancel_event.clear()
            if msg_type == "photo":
                photo_data = max(msg.get("photo", []), key=lambda p: p.get("file_size", 0))
                threading.Thread(
                    target=self._process_image_file,
                    args=(chat_id, photo_data, msg),
                    daemon=True,
                ).start()
            elif msg_type == "document":
                threading.Thread(
                    target=self._process_document_file,
                    args=(chat_id, msg.get("document"), msg),
                    daemon=True,
                ).start()
            elif msg_type == "voice":
                voice_data = msg.get("audio") or msg.get("voice")
                threading.Thread(
                    target=self._process_voice_file,
                    args=(chat_id, voice_data, msg),
                    daemon=True,
                ).start()
        else:
            pos = self._enqueue_message(chat_id, user_id, user_name, "", msg, msg_type=msg_type)
            type_label = {"photo": "圖片", "document": "文件", "voice": "語音"}.get(msg_type, msg_type)
            self.send(chat_id, f"⏳ {type_label}已排隊 #{pos} — 處理完當前任務後自動開始")

    def _buffer_media_group(self, chat_id: str, group_id: str, msg: dict):
        """Buffer a media group message. Start/reset a timer; when it fires, process all."""
        if chat_id not in self._media_group_buffers:
            self._media_group_buffers[chat_id] = {}
        buf = self._media_group_buffers[chat_id]
        buf.setdefault(group_id, []).append(msg)
        count = len(buf[group_id])

        # Cancel existing timer for this group, start a new one
        timer_key = f"{chat_id}:{group_id}"
        if timer_key in self._media_group_timers:
            self._media_group_timers[timer_key].cancel()

        timer = threading.Timer(self._media_group_wait, self._process_media_group, args=[chat_id, group_id])
        timer.daemon = True
        self._media_group_timers[timer_key] = timer
        timer.start()

        # Determine type label for the notification
        is_photo = bool(msg.get("photo"))
        label = "📸 圖片" if is_photo else "📎 檔案"
        self.send(chat_id, f"{label} {count} 收到 — 等埋其他同組檔案...")

    def _process_media_group(self, chat_id: str, group_id: str):
        """Timer fired — all media in this group should have arrived. Queue them all."""
        msgs = self._media_group_buffers.get(chat_id, {}).pop(group_id, [])
        timer_key = f"{chat_id}:{group_id}"
        self._media_group_timers.pop(timer_key, None)

        if not msgs:
            return

        total = len(msgs)
        is_photo = bool(msgs[0].get("photo"))
        label = "📸 圖片" if is_photo else "📎 檔案"
        self.send(chat_id, f"{label}組共 {total} 個 — 開始排隊處理...")

        for msg in msgs:
            user_id = str(msg["from"]["id"])
            user_name = msg["from"].get("first_name", "User")
            msg_type = "photo" if msg.get("photo") else "document"
            self._handle_media_msg(chat_id, user_id, user_name, msg, msg_type)

    # ─────────────────────────────────────────────────────────────
    #  File processors
    # ─────────────────────────────────────────────────────────────

    def _process_document_file(self, chat_id: str, doc: dict, msg: dict):
        """Download, extract, and analyze a document via BAW. Inline edit — one message."""
        try:
            file_id = doc["file_id"]
            file_name = doc.get("file_name", "document")

            # Pre-check: Telegram Bot API can't download files >20MB
            file_size = doc.get("file_size", 0)
            MAX_BOT_DOWNLOAD = 20 * 1024 * 1024  # 20MB
            if file_size > MAX_BOT_DOWNLOAD:
                size_mb = file_size / (1024 * 1024)
                self.send(chat_id,
                    f"📁 **{file_name}** ({size_mb:.1f}MB) 超過 Telegram Bot API 20MB 限制。\n"
                    f"請用其他方式傳送：壓縮、分割檔案、或改用雲端連結。"
                )
                return

            status_id = self.send(chat_id, f"📥 Downloading **{file_name}**...")

            local_path = self._download_file(file_id, file_name)
            self.send(chat_id, f"🔍 Extracting content...", edit_msg_id=status_id)

            content = self._extract_file_content(local_path)

            # ── Save extracted text to temp file so BAW's agent can use read_file/search_files ──
            # Instead of stuffing all text into the LLM prompt (which crashes on large PDFs),
            # we save the text and let the agent analyze it chunk-by-chunk using its tools.
            import tempfile as _tempfile
            _extracted_path = Path(_tempfile.gettempdir()) / f"baw_extracted_{file_name}.txt"
            _extracted_path.write_text(content, encoding="utf-8")
            _content_size = len(content)
            _content_pages = content.count("--- Page")  # rough page count

            prompt = (
                f"[File: {file_name}]\n"
                f"[Type: {doc.get('mime_type', 'unknown')}]\n\n"
                f"The full text of this document has been extracted and saved to:\n"
                f"  {_extracted_path}\n"
                f"  Size: {_content_size} chars, ~{_content_pages} pages\n\n"
                f"Use `read_file` with offset/limit to read it in chunks (500 lines at a time).\n"
                f"Use `search_files` with regex to find keywords across the document.\n"
                f"Do NOT try to read the entire file at once — read and analyze section by section.\n\n"
                f"Task: Analyze this file. Summarize its key content in Traditional Chinese. "
                f"If it's a technical document, identify the main topics. "
                f"Annotate key findings with page numbers from the extracted text."
            )

            self.send(chat_id, f"🤔 Analyzing **{file_name}** (~{_content_pages} pages)...", edit_msg_id=status_id)
            response = self._run_baw(prompt, chat_id=chat_id)
            try:
                self._client.post(
                    f"{self._api_base}/deleteMessage",
                    json={"chat_id": chat_id, "message_id": int(status_id)},
                    timeout=5,
                )
            except Exception:
                pass
            self.send(chat_id, response)
            self._record_batch_result(chat_id, response[:200], "document")

        except BaseException as e:
            if isinstance(e, (KeyboardInterrupt, SystemExit)):
                logger.warning(f"[Telegram] Document processing interrupted: {type(e).__name__}")
            else:
                logger.error(f"[Telegram] Document processing error: {e}")
            try:
                self.send(chat_id, f"❌ 無法處理文件: {e}"[:500])
            except Exception:
                pass
        finally:
            self._release_slot(chat_id)

    def _process_image_file(self, chat_id: str, photo_data: dict, msg: dict):
        """Download an image and analyze with vision AI (MiniMax API → Stepfun → OCR). Inline edit."""
        try:
            file_id = photo_data["file_id"]
            file_name = f"photo_{file_id[:8]}.jpg"
            status_id = self.send(chat_id, "📥 Downloading image...")
            local_path = self._download_file(file_id, file_name)

            from tools.vision import _vision_minimax, _vision_stepfun

            self.send(chat_id, "👁️ Analyzing with MiniMax vision...", edit_msg_id=status_id)

            # ── Primary: Direct MiniMax API (no CLI needed) ──
            vision_result = ""
            provider_used = ""
            try:
                vision_result = _vision_minimax(
                    local_path,
                    "Describe this image in detail. What objects, text, brands, or products do you see?"
                )
                if vision_result.startswith("Error:") or vision_result.startswith("MiniMax vision error:"):
                    raise RuntimeError(vision_result)
                provider_used = "MiniMax"
            except Exception as mmx_e:
                # ── Fallback: Stepfun multimodal ──
                self.send(chat_id, "👁️ MiniMax API failed, trying Stepfun...", edit_msg_id=status_id)
                try:
                    vision_result = _vision_stepfun(
                        local_path,
                        "Describe this image in detail. What objects, text, brands, or products do you see?"
                    )
                    if vision_result.startswith("Error:") or vision_result.startswith("Stepfun vision error:"):
                        raise RuntimeError(vision_result)
                    provider_used = "Stepfun"
                except Exception as step_e:
                    # ── Final fallback: OCR ──
                    self.send(chat_id, "📄 Falling back to OCR...", edit_msg_id=status_id)
                    content = self._extract_file_content(local_path)
                    vision_result = f"OCR: {content}"
                    provider_used = "OCR"

            prompt = (
                f"[Image analysis via {provider_used}]\n"
                f"File: {file_name}\n\n"
                f"Vision result:\n{vision_result}\n\n"
                f"---\n"
                f"Based on the vision analysis above, answer in Traditional Chinese:\n"
                f"- What is shown in this image?\n"
                f"- If it's a product: what is it, and where can I buy it?\n"
                f"- If there are similar items: suggest alternatives."
            )

            self.send(chat_id, "🤔 Analyzing with BAW...", edit_msg_id=status_id)
            response = self._run_baw(prompt, chat_id=chat_id)
            try:
                self._client.post(
                    f"{self._api_base}/deleteMessage",
                    json={"chat_id": chat_id, "message_id": int(status_id)},
                    timeout=5,
                )
            except Exception:
                pass
            self.send(chat_id, response)
            self._record_batch_result(chat_id, response[:200], "image")

        except Exception as e:
            logger.error(f"[Telegram] Image processing error: {e}")
            self.send(chat_id, f"❌ Error processing image: {e}")
        finally:
            self._release_slot(chat_id)

    def _process_voice_file(self, chat_id: str, voice_data: dict, msg: dict):
        """Download audio, check STT capability, then transcribe or present options."""
        import os
        text = ""
        info = None
        try:
            file_id = voice_data["file_id"]
            ext = ".ogg" if msg.get("voice") else ".mp3"
            self.send(chat_id, "📥 下載語音...")

            local_path = self._download_file(file_id, f"voice_{file_id[:8]}{ext}")

            # ── Step 1: Load BAW config and check STT capability ──
            baw = self._baw_ensure()
            config = baw["config"]
            default_model_id = config.get("model", {}).get("default", "")
            all_models = self._scan_all_models(config)

            # Current model capability
            current_model = next((m for m in all_models if m["id"] == default_model_id), None)
            audio_models = [m for m in all_models if m.get("audio_input")]

            # ── Step 2: Check STT capability ──
            stt_config = config.get("capabilities", {}).get("stt", {})
            stt_method = stt_config.get("method", "").strip()

            # Check if any model has "stt" capability
            stt_models = [m for m in all_models if "stt" in m.get("capabilities", [])]

            # ── Step 3: Check if faster-whisper is installed ──
            fw_available = False
            try:
                import faster_whisper as _fw
                fw_available = True
            except ImportError:
                pass

            # ── Capability message parts ──
            status_lines = []
            used_method = None

            # Priority A: Model native audio_input
            if current_model and current_model.get("audio_input"):
                status_lines.append(f"✅ 主模型 **{default_model_id}** 支援音訊輸入（audio_input=true）")
                # Native audio handling via model API — currently falls through
                # to configured STT methods or installed fallbacks below.
                # Future: implement native audio API call per model here.

            # Priority B: Config-defined STT method or model
            if not used_method and stt_method:
                status_lines.append(f"⚙️ STT 配置 method：**{stt_method}**")
                if stt_method == "faster-whisper" and fw_available:
                    fw_model = stt_config.get("model", "base")
                    status_lines.append(f"🎙️ 使用 faster-whisper（{fw_model}，本地免費）")
                    self.send(chat_id, "🎙️ 本地語音辨識中...")
                    try:
                        from faster_whisper import WhisperModel
                        logger.info(f"[Telegram] STT with faster-whisper ({fw_model}): {local_path}")
                        model = WhisperModel(fw_model, device="cpu", compute_type="int8")
                        segments, info = model.transcribe(local_path, language=None, beam_size=3)
                        text = " ".join(seg.text.strip() for seg in segments)
                        used_method = "faster-whisper"
                    except Exception as e:
                        logger.error(f"[Telegram] faster-whisper error: {e}")
                        self.send(chat_id, f"❌ faster-whisper 辨識失敗: {e}")
                        return
                elif stt_method == "faster-whisper" and not fw_available:
                    # Auto-install missing faster-whisper
                    self.send(chat_id, "⬇️ 正在安裝 faster-whisper 語音辨識引擎...")
                    import subprocess as _sp
                    _r = _sp.run(
                        [sys.executable, "-m", "pip", "install", "faster-whisper", "--quiet"],
                        capture_output=True, text=True, timeout=120,
                    )
                    if _r.returncode == 0:
                        fw_model = stt_config.get("model", "base")
                        self.send(chat_id, f"✅ 安裝成功。使用 faster-whisper（{fw_model}）...")
                        from faster_whisper import WhisperModel
                        model = WhisperModel(fw_model, device="cpu", compute_type="int8")
                        segments, info = model.transcribe(local_path, language=None, beam_size=3)
                        text = " ".join(seg.text.strip() for seg in segments)
                        used_method = "faster-whisper"
                elif stt_method == "openai-whisper":
                    api_key_env = stt_config.get("api_key_env", "OPENAI_API_KEY")
                    api_key = os.environ.get(api_key_env, "")
                    if api_key:
                        openai_model = stt_config.get("model", "whisper-1")
                        self.send(chat_id, "🎙️ OpenAI Whisper API 辨識中...")
                        try:
                            from openai import OpenAI
                            client = OpenAI(api_key=api_key)
                            with open(local_path, "rb") as f:
                                transcript = client.audio.transcriptions.create(
                                    model=openai_model, file=f
                                )
                            text = transcript.text or ""
                            used_method = "openai-whisper"
                        except Exception as e:
                            self.send(chat_id, f"❌ OpenAI Whisper API 失敗: {e}。嘗試其他方法...")
                    else:
                        status_lines.append(f"   ⚠️ {api_key_env} 未設定")
                elif stt_method in ("auto-asr", "model", "hybrid"):
                    # Auto-detect ASR protocol: try multiple endpoint patterns
                    # ── hybrid: resolve primary/fallback tiers before probing ──
                    if stt_method == "hybrid":
                        primary = stt_config.get("primary", {})
                        _active = primary if (primary.get("base_url") and os.environ.get(primary.get("api_key_env", ""))) else stt_config
                        api_key_env = _active.get("api_key_env", "")
                        stt_model_id = (_active.get("model") or stt_config.get("model", "")).strip()
                        base_url = _active.get("base_url", "")
                        if not base_url:
                            base_url = stt_config.get("base_url", "")
                        logger.info(f"[Telegram] hybrid STT: primary={'→'.join([primary.get(k,'') for k in ['model','base_url']]) if primary else 'none'}, resolved={stt_model_id} @ {base_url}")
                    else:
                        api_key_env = stt_config.get("api_key_env", stt_config.get("api_key_env", ""))
                        stt_model_id = stt_config.get("model", "").strip()
                        base_url = stt_config.get("base_url", "")
                    api_key = os.environ.get(api_key_env, "")
                    if api_key and base_url:
                        self.send(chat_id, f"🔍 ASR auto-detect（{base_url}）...")
                        try:
                            import httpx, base64 as _b64, json as _json
                            with open(local_path, "rb") as f:
                                audio_bytes = f.read()
                                audio_b64 = _b64.b64encode(audio_bytes).decode()
                            audio_type = "ogg" if local_path.endswith(".ogg") else "mp3"
                            b_url = base_url.rstrip("/")

                            # Strategy 1: OpenAI-compatible transcription endpoint
                            text = ""
                            for candidate_url in [b_url, b_url.replace("/v1", "/step_plan/v1")]:
                                try:
                                    transcription_url = f"{candidate_url}/audio/transcriptions"
                                    logger.info(f"[Telegram] ASR probe: OpenAI @ {transcription_url}")
                                    with httpx.Client(timeout=15, verify=True) as cli:
                                        files = {"file": (f"voice.{audio_type}", audio_bytes, f"audio/{audio_type}")}
                                        data = {"model": stt_model_id or "whisper-1", "language": "zh"}
                                        resp = cli.post(
                                            transcription_url, files=files, data=data,
                                            headers={"Authorization": f"Bearer {api_key}"},
                                        )
                                    if resp.status_code == 200:
                                        result = resp.json()
                                        text = result.get("text", "")
                                        if text:
                                            used_method = "openai-whisper"
                                            logger.info(f"[Telegram] ASR OK via OpenAI @ {candidate_url}")
                                            break
                                    elif resp.status_code == 402:
                                        logger.info(f"[Telegram] ASR 402 Quota @ {candidate_url}")
                                except Exception as e:
                                    logger.info(f"[Telegram] OpenAI-compatible ASR failed @ {candidate_url}: {e}")

                            # Strategy 1b: xAI/Grok-style STT endpoint (/v1/stt)
                            if not text:
                                try:
                                    stt_url = f"{b_url}/stt"
                                    logger.info(f"[Telegram] ASR probe: xAI/Grok @ {stt_url}")
                                    with httpx.Client(timeout=15, verify=True) as cli:
                                        files = {"file": (f"voice.{audio_type}", audio_bytes, f"audio/{audio_type}")}
                                        data = {"model": stt_model_id or "grok-stt", "language": "zh"}
                                        resp = cli.post(
                                            stt_url, files=files, data=data,
                                            headers={"Authorization": f"Bearer {api_key}"},
                                        )
                                    if resp.status_code == 200:
                                        result = resp.json()
                                        text = result.get("text", "") or result.get("content", "")
                                        if text:
                                            used_method = "openai-whisper"
                                            logger.info(f"[Telegram] ASR OK via xAI/Grok @ {stt_url}")
                                except Exception as e:
                                    logger.info(f"[Telegram] xAI/Grok STT failed @ {stt_url}: {e}")
                            # If we got 402, report it but still try SSE
                            if not text and resp and resp.status_code == 402:
                                self.send(chat_id, "⚠️ Stepfun API 配額不足（402），SSE 可能也一樣...")

                            # Strategy 2: Stepfun-style SSE ASR
                            if not text:
                                for sse_candidate in [b_url, b_url.replace("/v1", "/step_plan/v1")]:
                                    try:
                                        sse_url = f"{sse_candidate}/audio/asr/sse"
                                        logger.info(f"[Telegram] ASR probe: SSE @ {sse_url}")
                                        sse_payload = {
                                            "audio": {
                                                "data": audio_b64,
                                                "input": {
                                                    "transcription": {
                                                        "language": "zh",
                                                        "model": stt_model_id or "stepaudio-2.5-asr",
                                                        "enable_itn": True,
                                                        "enable_timestamp": False,
                                                    },
                                                    "format": {"type": audio_type},
                                                },
                                            }
                                        }
                                        full_text = ""
                                        with httpx.Client(timeout=60, verify=True) as cli:
                                            resp = cli.post(
                                                sse_url, json=sse_payload,
                                                headers={
                                                    "Authorization": f"Bearer {api_key}",
                                                    "Content-Type": "application/json",
                                                    "Accept": "text/event-stream",
                                                },
                                            )
                                            if resp.status_code in (200, 201):
                                                for line in resp.iter_lines():
                                                    line = line.strip()
                                                    if not line:
                                                        continue
                                                    if line.startswith("data: "):
                                                        data_str = line[6:]
                                                        if data_str == "[DONE]":
                                                            break
                                                        try:
                                                            evt = _json.loads(data_str)
                                                            t = evt.get("type", "")
                                                            if t == "transcript.text.delta":
                                                                full_text += evt.get("delta", "")
                                                            elif t == "transcript.text.done":
                                                                full_text = evt.get("text", full_text)
                                                                break
                                                            elif t == "error":
                                                                raise RuntimeError(evt.get("message", str(evt)))
                                                        except _json.JSONDecodeError:
                                                            continue
                                        if full_text.strip():
                                            text = full_text.strip()
                                            used_method = "auto-asr-sse"
                                            logger.info(f"[Telegram] ASR OK via SSE: {len(text)} chars")
                                    except Exception as e2:
                                        logger.info(f"[Telegram] SSE ASR failed: {e2}")

                            if text:
                                used_method = used_method or "auto-asr"
                            else:
                                # ── Hybrid fallback: primary failed, try fallback tier ──
                                if stt_method == "hybrid":
                                    fb = stt_config.get("fallback", {})
                                    if isinstance(fb, dict) and fb.get("base_url"):
                                        fb_api_env = fb.get("api_key_env", "")
                                        fb_key = os.environ.get(fb_api_env, "")
                                        fb_url = fb.get("base_url", "").rstrip("/")
                                        fb_model = fb.get("model", "")
                                        if fb_key and fb_url:
                                            logger.info(f"[Telegram] hybrid fallback: {fb_model} @ {fb_url}")
                                            # Strategy 1b retry with fallback endpoint
                                            try:
                                                stt_url = f"{fb_url}/stt"
                                                with httpx.Client(timeout=15, verify=True) as cli:
                                                    files = {"file": (f"voice.{audio_type}", audio_bytes, f"audio/{audio_type}")}
                                                    data = {"model": fb_model or "grok-stt", "language": "zh"}
                                                    resp = cli.post(
                                                        stt_url, files=files, data=data,
                                                        headers={"Authorization": f"Bearer {fb_key}"},
                                                    )
                                                if resp.status_code == 200:
                                                    result = resp.json()
                                                    text = result.get("text", "") or result.get("content", "")
                                                    if text:
                                                        used_method = "openai-whisper"
                                                        logger.info(f"[Telegram] hybrid fallback OK via {fb_url}/stt")
                                            except Exception as fe:
                                                logger.info(f"[Telegram] hybrid fallback failed: {fe}")
                                if not text:
                                    self.send(chat_id, "❌ 所有 ASR 協議都失敗（OpenAI + SSE），check API key/endpoint。")
                        except Exception as e:
                            logger.error(f"[Telegram] auto-asr error: {e}")
                            self.send(chat_id, f"❌ ASR auto-detect 失敗: {e}")
                elif stt_method == "google-speech":
                    status_lines.append("   ⚠️ Google Speech-to-Text 尚未實作")
                else:
                    status_lines.append(f"   ⚠️ 未知 STT method: {stt_method}")

            # Priority B2: Model-based STT (stt.model set but no stt.method)
            if not used_method:
                stt_model_id = stt_config.get("model", "").strip()
                if stt_model_id:
                    # Find the model in providers
                    stt_model_obj = next((m for m in all_models if m["id"] == stt_model_id), None)
                    if stt_model_obj and stt_model_obj.get("audio_input"):
                        # Use model's native audio input API
                        status_lines.append(f"🎙️ 使用模型原生音訊輸入：**{stt_model_id}**（audio_input=true）")
                        self.send(chat_id, f"🎙️ {stt_model_id} 音訊辨識中...")
                        try:
                            import base64 as _b64
                            provider = stt_model_obj["provider"]
                            # Find provider config for API credentials
                            pconf = config.get("providers", {}).get(provider, {})
                            api_key = os.environ.get(pconf.get("api_key_env", ""), "")
                            api_url = f"{pconf.get('base_url', '').rstrip('/')}/chat/completions"
                            with open(local_path, "rb") as f:
                                audio_b64 = _b64.b64encode(f.read()).decode()
                            data_url = f"data:audio/ogg;base64,{audio_b64}"
                            resp = httpx.post(
                                api_url,
                                json={
                                    "model": stt_model_id,
                                    "messages": [{
                                        "role": "user",
                                        "content": [
                                            {"type": "text", "text": "Transcribe this audio to text in Traditional Chinese (Cantonese). Return ONLY the transcription, no extra text."},
                                            {"type": "audio_url", "audio_url": {"url": data_url}},
                                        ],
                                    }],
                                    "max_tokens": 1024,
                                },
                                headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
                                timeout=60,
                            )
                            data = resp.json()
                            text = (data.get("choices", [{}])[0]
                                    .get("message", {})
                                    .get("content", ""))
                            used_method = f"model-{stt_model_id}"
                        except Exception as e:
                            logger.error(f"[Telegram] Model STT error: {e}")
                            self.send(chat_id, f"❌ {stt_model_id} 音訊辨識失敗: {e}")
                    else:
                        status_lines.append(f"   ⚠️ stt.model={stt_model_id} 無 audio_input 能力")

            # Priority C: faster-whisper local fallback (only if no STT config at all)
            if not used_method and fw_available and not stt_config.get("model", "").strip():
                fw_model = stt_config.get("model", "base") if stt_method == "faster-whisper" else "base"
                status_lines.append(f"🎙️ 使用 faster-whisper（{fw_model}，本地免費）")
                self.send(chat_id, "🎙️ 本地語音辨識中...")
                try:
                    from faster_whisper import WhisperModel
                    logger.info(f"[Telegram] STT with faster-whisper ({fw_model}): {local_path}")
                    model = WhisperModel(fw_model, device="cpu", compute_type="int8")
                    segments, info = model.transcribe(local_path, language=None, beam_size=3)
                    text = " ".join(seg.text.strip() for seg in segments)
                    used_method = "faster-whisper"
                except Exception as e:
                    logger.error(f"[Telegram] faster-whisper error: {e}")
                    self.send(chat_id, f"❌ faster-whisper 辨識失敗: {e}")
                    return

            # ── Result: transcribe success or present options ──
            if used_method:
                if not text.strip():
                    self.send(chat_id, "🔇 聽唔到任何語音內容。")
                    return

                dur_str = "?"
                if used_method == "faster-whisper" and info and hasattr(info, 'duration') and isinstance(info.duration, (int, float)):
                    dur_str = f"{info.duration:.1f}s"
                self.send(chat_id, f"📝 辨識完成（{dur_str}）：\n\n{text[:200]}{'…' if len(text)>200 else ''}")

                prompt = (
                    f"[語音輸入 — 用戶語音訊息，已自動轉文字]\n"
                    f"用戶說：{text}\n\n"
                    f"請用繁體中文回應，分析並回答對方提問。"
                )
                self.send(chat_id, "🤔 BAW 分析中...")
                response = self._run_baw(prompt, chat_id=chat_id)
                self.send(chat_id, response)
                self._record_batch_result(chat_id, response[:200], "voice")
            else:
                # ── No STT method available — present diagnostics + options ──
                stt_method = stt_config.get("method", "")
                # hybrid: resolve primary config for accurate diagnostics
                if stt_method == "hybrid":
                    primary = stt_config.get("primary", {})
                    if primary.get("base_url") and os.environ.get(primary.get("api_key_env", "")):
                        stt_model = primary.get("model", stt_config.get("model", ""))
                        stt_base = primary.get("base_url", "")
                        stt_key_env = primary.get("api_key_env", "")
                    else:
                        stt_model = stt_config.get("model", "")
                        stt_base = stt_config.get("base_url", "")
                        stt_key_env = stt_config.get("api_key_env", "")
                else:
                    stt_model = stt_config.get("model", "")
                    stt_base = stt_config.get("base_url", "")
                    stt_key_env = stt_config.get("api_key_env", "")
                key_set = bool(os.environ.get(stt_key_env)) if stt_key_env else False

                if stt_method and stt_model:
                    # STT IS configured — explain why it failed
                    reason = ""
                    if stt_key_env and not key_set:
                        reason = f"`{stt_key_env}` 未 set 喺 `.env`"
                    elif stt_base and stt_key_env:
                        reason = f"`{stt_method}` probe `{stt_base}` 失敗（API key 有效但端點無回應）"
                    else:
                        reason = f"`{stt_method}` probe 失敗"
                    msg_parts = [
                        f"🎵 收到語音訊息\\n\\n"
                        f"STT 已設定：`{stt_method}` → `{stt_model}` @ `{stt_base or 'auto'}`\\n"
                        f"失敗原因：{reason}\\n"
                    ]
                else:
                    msg_parts = ["🎵 收到語音訊息，但目前 BAW 未有語音處理能力。\\n"]

                # If faster-whisper is installed and viable, mention it briefly
                if fw_available and not used_method:
                    msg_parts.append(
                        "💡 faster-whisper 已安裝，改 `stt.method: faster-whisper` 可即用（本地免費）"
                    )

                self.send(chat_id, "\n".join(msg_parts))

        except Exception as e:
            logger.error(f"[Telegram] Voice processing error: {e}")
            self.send(chat_id, f"❌ 語音處理錯誤: {e}")
        finally:
            self._release_slot(chat_id)

    @staticmethod
    def _scan_all_models(config: dict) -> list[dict]:
        """Scan all providers/models from config and return flat list with capability info."""
        models = []
        providers = config.get("providers", {})
        for pname, pconf in providers.items():
            for m in pconf.get("models", []):
                models.append({
                    "id": m.get("id", ""),
                    "provider": pname,
                    "vision": m.get("vision", False),
                    "audio_input": m.get("audio_input", False),
                    "capabilities": m.get("capabilities", []),
                })
        return models

    def _send_as_tts(self, chat_id: str, text: str):
        """Convert response text to speech and send as audio message."""
        import os
        try:
            # Load BAW config to get TTS credentials
            baw = self._baw_ensure()
            config = baw["config"]

            # Find TTS capability
            from core.capabilities import resolve_capability
            tts_cap = resolve_capability(config, "tts")
            if tts_cap is None or tts_cap["type"] != "model":
                logger.warning("[TTS] No TTS-capable model configured")
                return

            # Get API key
            api_key = os.environ.get(tts_cap.get("api_key_env", ""), "")
            if not api_key:
                logger.warning(f"[TTS] No API key for {tts_cap.get('id', '?')}")
                return

            # Get TTS config
            tts_config = config.get("capabilities", {}).get("tts", {}).get("config", {})
            tts_model = tts_config.get("api_model", "speech-2.8-hd")
            voice = self._tts_voice or tts_config.get("voice", "male-tone-1")

            # Call TTS
            from core.tts import minimax_tts, save_audio_bytes
            audio_data = minimax_tts(
                text=text,
                api_key=api_key,
                voice=voice,
                model=tts_model,
                language="yue",
            )
            if audio_data is None:
                self.send(chat_id, "⚠️ TTS 生成失敗")
                return

            audio_path = save_audio_bytes(audio_data)
            self.send(chat_id, f"MEDIA:{audio_path}")

        except Exception as e:
            logger.error(f"[TTS] send error: {e}")

    def _poll_loop(self):
        """Long-polling loop."""
        if not self._client:
            logger.error("[Telegram] No client — can't poll")
            return

        while self._running:
            try:
                r = self._client.post(
                    f"{self._api_base}/getUpdates",
                    json={
                        "offset": self._offset,
                        "timeout": POLL_TIMEOUT,
                        "allowed_updates": ["message", "callback_query"],
                    },
                    timeout=POLL_TIMEOUT + 5,
                )
                if r.status_code != 200:
                    logger.warning(f"[Telegram] getUpdates HTTP {r.status_code}")
                    time.sleep(POLL_RETRY_DELAY)
                    continue

                data = r.json()
                if not data.get("ok"):
                    logger.warning(f"[Telegram] API error: {data}")
                    time.sleep(POLL_RETRY_DELAY)
                    continue

                for update in data.get("result", []):
                    self._offset = update["update_id"] + 1
                    self._handle_update(update)

            except httpx.TimeoutException:
                # Timeout is normal for long-poll — just retry
                continue
            except Exception as e:
                logger.error(f"[Telegram] Poll error: {e}")
                time.sleep(POLL_RETRY_DELAY)

    def _handle_update(self, update: dict):
        """Process a single Telegram update (non-blocking, concurrent)."""
        # ── Handle inline keyboard callback query ──
        cb = update.get("callback_query")
        if cb:
            self._handle_callback(cb)
            return

        msg = update.get("message")
        if not msg:
            return

        chat_id = str(msg["chat"]["id"])
        user_id = str(msg["from"]["id"])
        user_name = msg["from"].get("first_name", "User")
        text = msg.get("text", "").strip()

        # ── Handle reply-to: prepend quoted message context
        reply = msg.get("reply_to_message")
        if reply and text:
            reply_text = reply.get("text", "") or reply.get("caption", "") or ""
            if reply_text:
                reply_from = reply.get("from", {}).get("first_name", "User")
                text = f"> {reply_from}: {reply_text[:200]}\n\n{text}"

        if not text:
            # ── Non-text message — check media group FIRST ──
            media_group_id = msg.get("media_group_id")
            if media_group_id:
                self._buffer_media_group(chat_id, media_group_id, msg)
                return

            # ── Single non-text message (no group) — queue if busy, else process ──
            doc = msg.get("document")
            photo = msg.get("photo")
            video = msg.get("video")
            audio = msg.get("audio")
            voice = msg.get("voice")

            if doc:
                self._handle_media_msg(chat_id, user_id, user_name, msg, "document")
            elif photo:
                self._handle_media_msg(chat_id, user_id, user_name, msg, "photo")
            elif video:
                self.send(chat_id, "🎬 收到影片，但 BAW 暫時未支援影片處理。")
            elif audio or voice:
                self._handle_media_msg(chat_id, user_id, user_name, msg, "voice")
            # else: silent ignore for unknown types
            return

        # Access control
        if self._allowed and user_id not in self._allowed:
            self.send(chat_id, "⛔ You are not authorized to use this bot.")
            return

        logger.info(f"[Telegram] <{user_name}> {text[:80]}")

        # /stop — cancel running BAW immediately
        if text.strip().lower().startswith("/stop"):
            self._cancel_event.set()
            self.send(chat_id, f"⏹ Stopped {self._active_count} active task(s).")
            return

        # /tts — toggle text-to-speech
        if text.strip().lower().startswith("/tts"):
            args = text.strip().lower().split()
            if len(args) > 1:
                if args[1] in ("on", "true", "1", "yes"):
                    self._tts_enabled = True
                    self.send(chat_id, "🔊 TTS 已開啟 — 回覆會自動轉語音")
                elif args[1] in ("off", "false", "0", "no"):
                    self._tts_enabled = False
                    self.send(chat_id, "🔇 TTS 已關閉")
                else:
                    self._tts_voice = args[1]
                    self.send(chat_id, f"🎤 TTS 語音設為: {args[1]}")
            else:
                status = "開" if self._tts_enabled else "關"
                self.send(chat_id,
                          f"🔊 TTS 狀態: **{status}**\n"
                          f"   語音: `{self._tts_voice}`\n"
                          f"   用 `/tts on` 開啟\n"
                          f"   用 `/tts off` 關閉\n"
                          f"   用 `/tts <voice_id>` 切換語音")
            return

        # /selftest — run system diagnostics
        if text.strip().lower().startswith("/selftest"):
            self.send(chat_id, "🧪 Running self-test...")
            try:
                from tools.selftest import selftest as _st
                report = _st(full=False)
                self.send(chat_id, report)
            except Exception as e:
                self.send(chat_id, f"❌ Self-test failed: {e}")
            return

        # ── /btw bypass: quick answer, skip queue entirely ──
        _btw_match = text.lstrip().lower().startswith("/btw")
        if _btw_match:
            _btw_text = text.lstrip()[4:].strip()
            if not _btw_text:
                self.send(chat_id, "/btw `<question>` — Quick answer, no queue")
                return
            threading.Thread(
                target=self._process_btw,
                args=(chat_id, _btw_text),
                daemon=True,
            ).start()
            return

        # Debounce window: suppress new threads right after /stop
        if self._debounce_until and time.time() < self._debounce_until:
            self.send(chat_id, "⏳ Please wait a moment before sending a new request.")
            return

        # Acquire slot and start async processing in background thread
        if not self._acquire_slot():
            pos = self._enqueue_message(chat_id, user_id, user_name, text, msg)
            self.send(chat_id, f"⏳ Queued #{pos} — will process when a slot frees up ({self._active_count}/{self._max_concurrency} busy)")
            return

        # Per-chat sequential: if this chat already has a task, queue it
        if self._is_chat_busy(chat_id):
            with self._active_lock:
                self._active_count = max(0, self._active_count - 1)
            pos = self._enqueue_message(chat_id, user_id, user_name, text, msg)
            self.send(chat_id, f"⏳ Queued #{pos} — another task is running in this chat")
            return
        self._mark_chat_busy(chat_id)

        self._cancel_event.clear()
        threading.Thread(
            target=self._process_message,
            args=(chat_id, user_id, user_name, text, msg),
            daemon=True,
        ).start()

    def _set_reaction(self, chat_id: int, message_id: int, emoji: str) -> bool:
        """Set a reaction emoji on a telegram message. Falls back to big emoji if unsupported."""
        try:
            resp = self._client.post(
                f"{self._api_base}/setMessageReaction",
                json={
                    "chat_id": chat_id,
                    "message_id": message_id,
                    "reaction": [{"type": "emoji", "emoji": emoji}],
                },
                timeout=5,
            )
            data = resp.json()
            if data.get("ok"):
                return True
            # Fallback: big reaction emoji
            fallback = {"🔄": "⚡", "🧠": "🤔", "✅": "👍", "❌": "👎"}
            fb = fallback.get(emoji)
            if fb:
                return self._set_reaction(chat_id, message_id, fb)
            return False
        except Exception:
            return False

    def _process_btw(self, chat_id: str, question: str):
        """Quick LLM answer — no agent loop, no court, no queue. Runs in own thread."""
        try:
            from ..llm import call_llm_with_fallback, calculate_cost, get_model
            from ..loop import build_system_prompt
            config = self._baw_config
            data_dir = self._data_dir
            model = get_model(config)
            # Minimal system prompt — no soul, no memories, no rules
            sys_prompt = build_system_prompt(config, data_dir, fresh_start=True)
            msgs = [
                {"role": "system", "content": f"你係 BAW。快速回答問題。\n\n{sys_prompt}"},
                {"role": "user", "content": question},
            ]
            fb = call_llm_with_fallback(config, msgs, temperature=0.7)
            answer = fb.response.content.strip() if fb.response else "(no response)"
            cost = calculate_cost(model, fb.response.input_tokens, fb.response.output_tokens) if fb.response else 0
            total_tokens = (fb.response.input_tokens or 0) + (fb.response.output_tokens or 0) if fb.response else 0
            self.send(chat_id,
                      f"⚡ <b>BTW</b>: {answer}\n"
                      f"📊 {cost:.4f} | {total_tokens:,} tokens")
        except Exception as e:
            self.send(chat_id, f"❌ BTW error: {e}")

    def _process_message(self, chat_id, user_id, user_name, text, msg):
        """Process a message in background thread (runs route() + send())."""
        logger.info(f"[Telegram] Processing: {text[:60]}...")

        # Get message_id for reactions
        _msg_id = msg.get("message_id") if isinstance(msg, dict) else None

        # Set initial reaction: 🔄 working/spinning
        if _msg_id:
            self._set_reaction(chat_id, _msg_id, "🔄")

        # Reaction update: after 2s switch to thinking
        _reaction_stop = threading.Event()
        if _msg_id:
            def _reaction_delayed():
                _reaction_stop.wait(2.0)
                if not _reaction_stop.is_set():
                    self._set_reaction(chat_id, _msg_id, "🧠")
            _rh = threading.Thread(target=_reaction_delayed, daemon=True)
            _rh.start()

        # M2 (Fable 5 spec): send a placeholder court message immediately, then
        # edit-in-place as the case progresses. This gives the user instant
        # "filed" feedback (TTFT < 500ms) instead of waiting silently for the
        # full verdict. The placeholder is a single ⚖️ emoji + the user's
        # question, which is overwritten in the success path below.
        _placeholder_msg_id = ""
        if not text.lstrip().startswith("/"):  # don't placeholder slash commands
            _placeholder_msg_id = self.send(
                chat_id,
                f"🔧 {text[:50]}{'…' if len(text) > 50 else ''}",
            )

        # ── Typing indicator heartbeat ──
        _typing_stop = threading.Event()
        _typing_count = [0]  # mutable counter for debug
        _typing_start = [time.time()]  # mutable timestamp for elapsed calc

        def _typing_heartbeat():
            while not _typing_stop.is_set() and not self._cancel_event.is_set():
                for _retry in range(3):
                    try:
                        self._client.post(
                            f"{self._api_base}/sendChatAction",
                            json={"chat_id": chat_id, "action": "typing"},
                            timeout=5,
                        )
                        _typing_count[0] += 1
                        break
                    except Exception as _te:
                        logger.debug(
                            f"[Telegram] typing heartbeat retry {_retry + 1} "
                            f"(#{_typing_count[0]}): {_te}"
                        )
                        _typing_stop.wait(0.5)
                _typing_stop.wait(3.0)
        _hb = threading.Thread(
            target=_typing_heartbeat, daemon=True, name=f"typing-hb-{chat_id}",
        )
        _hb.start()

        # ── Progress monitor: send periodic "still working" updates ──
        _progress_stop = threading.Event()

        def _progress_monitor():
            _milestones = [120, 420, 720, 1200, 1800, 3000]  # 2min, 7min, 12min, 20min, 30min, 50min
            _msg_sent = False
            while not _progress_stop.wait(1.0) and not _typing_stop.is_set() and not self._cancel_event.is_set():
                _elapsed = int(time.time() - _typing_start[0])
                for _m in _milestones:
                    if _elapsed >= _m and not getattr(_progress_monitor, f"_done_{_m}", False):
                        setattr(_progress_monitor, f"_done_{_m}", True)
                        _msg_sent = True
                        if _elapsed < 600:
                            self.send(chat_id, f"⏳ Still working on your request... ({_elapsed // 60} min)")
                        else:
                            self.send(chat_id, f"⏳ Still working... ({_elapsed // 60} min). If stuck, use /stop to cancel.")
                        break
            # Final message if we ever sent progress updates
            if _msg_sent and not self._cancel_event.is_set():
                _elapsed_final = int(time.time() - _typing_start[0])
                _min_final = _elapsed_final // 60
                _sec_final = _elapsed_final % 60
                if _min_final > 0:
                    self.send(chat_id, f"Done ({_min_final} min {_sec_final}s)")
                else:
                    self.send(chat_id, f"Done ({_sec_final}s)")
        _pm = threading.Thread(
            target=_progress_monitor, daemon=True, name=f"progress-{chat_id}",
        )
        _pm.start()

        try:
            # Check cancelled before starting BAW
            if self._cancel_event.is_set():
                return

            msg_obj = Message(
                platform="telegram",
                chat_id=chat_id,
                user_id=user_id,
                text=text,
                raw=msg,
            )
            logger.info(
                f"[Telegram] typing heartbeat started "
                f"(chat_id={chat_id}, message='{text[:40]}…')"
            )
            response = self.route(msg_obj)
            _typing_stop.set()
            _reaction_stop.set()
            logger.info(
                f"[Telegram] typing heartbeat stopped "
                f"(chat_id={chat_id}, fired {_typing_count[0]} times, "
                f"elapsed={time.time() - _typing_start[0]:.1f}s)"
            )

            # Set reaction: ✅ success
            if _msg_id:
                self._set_reaction(chat_id, _msg_id, "✅")

            # If cancelled during processing, discard result
            if self._cancel_event.is_set():
                _reaction_stop.set()
                return
            if not response or not response.strip():
                logger.warning(f"[Telegram] Empty response from BAW for: {text[:80]}")
                response = "已完成。如果你看不到預期結果，可能是因為：\n(1) 該操作沒有輸出 (如寫檔、設定配置)\n(2) 系統發生錯誤但未回傳\n\n請嘗試 `/status` 或 `/doctor` 檢查狀態。"
            if response:
                logger.info(f"[Telegram] Sending response (len={len(response)}, placeholder={bool(_placeholder_msg_id)})")
                # M2: edit the placeholder in-place if we sent one, else send fresh.
                if _placeholder_msg_id:
                    _edit_ok = self.send(chat_id, response, edit_msg_id=_placeholder_msg_id)
                    if not _edit_ok:
                        # Edit failed — send as new message so user still sees the result
                        logger.warning(f"[Telegram] Placeholder edit failed, sending new message")
                        self.send(chat_id, response)
                else:
                    self.send(chat_id, response)
                self._record_batch_result(chat_id, response[:200], "text")
                # TTS: convert response to speech if enabled
                if self._tts_enabled and response.strip():
                    self._send_as_tts(chat_id, response)

            # If restart was requested, exit cleanly so systemd restarts
            if self._restart_requested:
                import os
                logger.info("[Telegram] Restart requested — exiting")
                os._exit(0)
        except Exception as e:
            logger.error(f"[Telegram] Process error: {e}")
            _typing_stop.set()
            _reaction_stop.set()
            # Set reaction: ❌ error
            if _msg_id:
                self._set_reaction(chat_id, _msg_id, "❌")
            if not self._cancel_event.is_set():
                try:
                    _err_str = str(e)
                    # If error already has a friendly emoji header, use it directly
                    if _err_str.startswith(("🚫", "⚠️", "✅", "ℹ️")):
                        err_text = _err_str
                    else:
                        err_text = f"❌ {_err_str[:800]}"
                        if len(_err_str) > 800:
                            err_text += "\n\n(錯誤訊息已截斷, 完整記錄在 log 中)"
                    if _placeholder_msg_id:
                        self.send(chat_id, err_text, edit_msg_id=_placeholder_msg_id)
                    else:
                        self.send(chat_id, err_text)
                except Exception:
                    pass
        finally:
            self._release_slot(chat_id)
    # ── Inline keyboard for model selection (hierarchical: provider → model) ──

    def _get_providers_config(self) -> dict:
        """Get provider → models mapping from config + auto-discovered models."""
        try:
            baw = self._baw_ensure()
            config = baw["config"]
            providers = config.get("providers", {})
            result = {}
            for pname, pcfg in providers.items():

                models = [m["id"] for m in pcfg.get("models", [])]
                result[pname] = models
            # Try auto-discovery: fetch /v1/models for each provider
            discovered = self._discover_provider_models(providers)
            for pname, extra in discovered.items():
                existing = set(result.get(pname, []))
                new = [m for m in extra if m not in existing]
                if new:
                    result.setdefault(pname, []).extend(new)
            return result
        except Exception:
            # Fallback hardcoded
            return {
                "deepseek": ["deepseek-v4-flash"],
                "kimi": ["kimi-k2.6"],
                "minimax": ["MiniMax-M2.5"],
            }

    def _discover_provider_models(self, providers: dict) -> dict:
        """Try to auto-discover available models from each provider's API."""
        import os
        discovered = {}
        for pname, pcfg in providers.items():
            base_url = pcfg.get("base_url", "").rstrip("/")
            api_key = os.environ.get(pcfg.get("api_key_env", ""), "")
            if not base_url or not api_key:
                continue
            models_url = f"{base_url}/models"
            try:
                resp = httpx.get(
                    models_url,
                    headers={"Authorization": f"Bearer {api_key}"},
                    timeout=5,
                )
                if resp.status_code != 200:
                    continue
                data = resp.json()
                # OpenAI-compatible format: {"data": [{"id": "..."}]}
                raw = data.get("data", [])
                if raw and isinstance(raw, list):
                    ids = [m.get("id", "") for m in raw if m.get("id")]
                    if ids:
                        discovered[pname] = ids
            except Exception:
                continue
        return discovered

    def _get_current_model_for_role(self, role: str) -> str | None:
        """Read current model for a role from config. Returns None if unset (inherits default)."""
        import yaml
        from pathlib import Path
        cfg_path = Path.home() / ".baw" / "config.yaml"
        if not cfg_path.exists():
            return None
        cfg = yaml.safe_load(cfg_path.read_text(encoding="utf-8"))
        model_cfg = cfg.get("model", {})
        adv_cfg = cfg.get("adversarial", {})
        exec_cfg = cfg.get("executor", {})
        key_map = {
            "default": model_cfg.get("default"),
            "angel": adv_cfg.get("angel_model"),
            "devil": adv_cfg.get("devil_model"),
            "executor": exec_cfg.get("model"),
        }
        return key_map.get(role)

    def _send_role_selector(self, chat_id: str, text: str) -> bool:
        """Show role selection keyboard with current settings."""
        import yaml
        from pathlib import Path
        cfg_path = Path.home() / ".baw" / "config.yaml"
        default_model = "unknown"
        if cfg_path.exists():
            cfg = yaml.safe_load(cfg_path.read_text(encoding="utf-8"))
            default_model = cfg.get("model", {}).get("default", "unknown")

        role_order = [("default", "Default Model"), ("angel", "Angel Model"), ("devil", "Devil Model"), ("executor", "Executor Model")]
        rows = []
        status_lines = []
        for role_key, label in role_order:
            current = self._get_current_model_for_role(role_key)
            if current and role_key != "default":
                rows.append([{"text": f"  {label}: {current} ✎", "callback_data": f"role_select:{role_key}"}])
                status_lines.append(f"  {label}: `{current}`")
            elif current and role_key == "default":
                rows.append([{"text": f"  {label}: {current}", "callback_data": f"role_select:{role_key}"}])
                status_lines.append(f"  {label}: `{current}`")
            else:
                rows.append([{"text": f"  {label} = Default ({default_model})", "callback_data": f"role_select:{role_key}"}])
                status_lines.append(f"  {label}: = Default (`{default_model}`)")
        try:
            resp = self._client.post(
                f"{self._api_base}/sendMessage",
                json={
                    "chat_id": chat_id,
                    "text": "<b>Model Configuration</b>\n" + "\n".join(status_lines),
                    "reply_markup": {"inline_keyboard": rows},
                    "parse_mode": "HTML",
                },
                timeout=5,
            )
            data = resp.json()
            if data.get("ok"):
                self._selector_chat = chat_id
                self._selector_msg_id = data["result"]["message_id"]
            return data.get("ok", False)
        except Exception as e:
            logger.warning(f"[Telegram] Role selector error: {e}")
            return False

    def _send_role_buttons(self, chat_id: str, msg_id: int):
        """Edit existing message to show role selection buttons with current settings."""
        import yaml
        from pathlib import Path
        cfg_path = Path.home() / ".baw" / "config.yaml"
        default_model = "unknown"
        if cfg_path.exists():
            cfg = yaml.safe_load(cfg_path.read_text(encoding="utf-8"))
            default_model = cfg.get("model", {}).get("default", "unknown")

        role_order = [("default", "Default Model"), ("angel", "Angel Model"), ("devil", "Devil Model"), ("executor", "Executor Model")]
        rows = []
        status_lines = []
        for role_key, label in role_order:
            current = self._get_current_model_for_role(role_key)
            if current and role_key != "default":
                rows.append([{"text": f"  {label}: {current} ✎", "callback_data": f"role_select:{role_key}"}])
                status_lines.append(f"  {label}: `{current}`")
            elif current and role_key == "default":
                rows.append([{"text": f"  {label}: {current}", "callback_data": f"role_select:{role_key}"}])
                status_lines.append(f"  {label}: `{current}`")
            else:
                rows.append([{"text": f"  {label} = Default ({default_model})", "callback_data": f"role_select:{role_key}"}])
                status_lines.append(f"  {label}: = Default (`{default_model}`)")
        try:
            self._client.post(
                f"{self._api_base}/editMessageText",
                json={
                    "chat_id": chat_id,
                    "message_id": msg_id,
                    "text": "<b>Model Configuration</b>\n" + "\n".join(status_lines),
                    "reply_markup": {"inline_keyboard": rows},
                    "parse_mode": "HTML",
                },
                timeout=5,
            )
        except Exception as e:
            logger.warning(f"[Telegram] Role buttons edit error: {e}")

    def _send_model_selector_text(self, chat_id: str, text: str) -> bool:
        """Parse [MODEL_SELECT] formatted text and send provider-level keyboard."""
        lines = text.strip().split("\n")
        title = lines[1] if len(lines) > 1 else "<b>Select Provider</b>"
        current_model = lines[2] if len(lines) > 2 else ""
        # Build provider keyboard
        providers = self._get_providers_config()
        rows = []
        for pname in providers:
            mark = "●" if any(current_model == m for m in providers[pname]) else " "
            rows.append([{"text": f"{mark} {pname}", "callback_data": f"provider_select:{pname}"}])
        if not rows:
            return self._send_text(chat_id, title)
        try:
            resp = self._client.post(
                f"{self._api_base}/sendMessage",
                json={
                    "chat_id": chat_id,
                    "text": f"{title}\nCurrent: <code>{current_model}</code>",
                    "reply_markup": {"inline_keyboard": rows},
                    "parse_mode": "HTML",
                },
                timeout=5,
            )
            data = resp.json()
            if data.get("ok"):
                self._selector_chat = chat_id
                self._selector_msg_id = data["result"]["message_id"]
            return data.get("ok", False)
        except Exception as e:
            logger.warning(f"[Telegram] Provider selector send error: {e}")
            return False

    def _send_model_buttons(self, chat_id: str, msg_id: int, provider: str, current_model: str):
        """Edit message to show models for a specific provider."""
        providers = self._get_providers_config()
        models = providers.get(provider, [])
        if not models:
            return
        rows = []
        for m in models:
            label = f"\u25cf {m}" if m == current_model else f"  {m}"
            rows.append([{"text": label, "callback_data": f"model_select:{m}"}])
        # Back button
        rows.append([{"text": "\u2190 Back", "callback_data": "provider_select:__back__"}])
        try:
            self._client.post(
                f"{self._api_base}/editMessageText",
                json={
                    "chat_id": chat_id,
                    "message_id": msg_id,
                    "text": f"<b>{provider}</b> models:\nCurrent: <code>{current_model}</code>",
                    "reply_markup": {"inline_keyboard": rows},
                    "parse_mode": "HTML",
                },
                timeout=5,
            )
        except Exception as e:
            logger.warning(f"[Telegram] Model buttons error: {e}")

    def _handle_callback(self, cb: dict):
        """Handle inline keyboard callback (role → provider → model selection)."""
        data = cb.get("data", "")
        msg = cb.get("message", {})
        chat_id = str(msg["chat"]["id"])
        cb_id = cb["id"]
        msg_id = msg.get("message_id")
        logger.info(f"[Telegram] Callback: {data} from {chat_id}")

        # ── Critical: acknowledge callback FIRST to stop Telegram loading spinner ──
        self._client.post(
            f"{self._api_base}/answerCallbackQuery",
            json={"callback_query_id": cb_id},
            timeout=5,
        )

        if data.startswith("role_select:"):
            role = data.split(":", 1)[1]
            if role == "__back__":
                # Go back to role selector (called from provider screen)
                self._send_role_buttons(chat_id, msg_id)
                return
            self._selector_role[chat_id] = role
            role_names = {"default": "Default", "angel": "Angel", "devil": "Devil", "executor": "Executor"}
            role_name = role_names.get(role, role)

            # Show provider keyboard for this role
            providers = self._get_providers_config()
            if msg_id:
                rows = []
                for pname in providers:
                    rows.append([{"text": pname, "callback_data": f"provider_select:{pname}"}])
                rows.append([{"text": "← Back", "callback_data": "role_select:__back__"}])
                try:
                    self._client.post(
                        f"{self._api_base}/editMessageText",
                        json={
                            "chat_id": chat_id,
                            "message_id": msg_id,
                            "text": f"<b>{role_name} Model</b>\nSelect provider:",
                            "reply_markup": {"inline_keyboard": rows},
                            "parse_mode": "HTML",
                        },
                        timeout=5,
                    )
                except Exception as e:
                    logger.warning(f"[Telegram] Role→provider error: {e}")

        elif data.startswith("provider_select:"):
            pname = data.split(":", 1)[1]
            if pname == "__back__":
                # If we came from a role selector, go back there; else go to provider list
                role = self._selector_role.get(chat_id, "")
                if role and msg_id:
                    self._send_role_buttons(chat_id, msg_id)
                else:
                    try:
                        providers = self._get_providers_config()
                        title = "<b>Select Provider</b>"
                        cc = getattr(self, '_chat_config', {}).get(chat_id, {})
                        current = cc.get("model", "deepseek-v4-flash")
                        rows = []
                        for pn in providers:
                            rows.append([{"text": f"  {pn}", "callback_data": f"provider_select:{pn}"}])
                        self._client.post(
                            f"{self._api_base}/editMessageText",
                            json={
                                "chat_id": chat_id,
                                "message_id": msg_id,
                                "text": f"{title}\nCurrent: <code>{current}</code>",
                                "reply_markup": {"inline_keyboard": rows},
                                "parse_mode": "HTML",
                            },
                            timeout=5,
                        )
                    except Exception as e:
                        logger.warning(f"[Telegram] Back button error: {e}")
            else:
                # Show models for this provider
                providers = self._get_providers_config()
                cc = getattr(self, '_chat_config', {}).get(chat_id, {})
                current = cc.get("model", "deepseek-v4-flash")
                self._send_model_buttons(chat_id, msg_id, pname, current)

        elif data.startswith("model_select:"):
            model_name = data.split(":", 1)[1]
            role = self._selector_role.get(chat_id, "default")

            # Map role to config key
            config_key_map = {
                "default": "model.default",
                "angel": "adversarial.angel_model",
                "devil": "adversarial.devil_model",
                "executor": "executor.model",
            }
            config_key = config_key_map.get(role, "model.default")
            role_names = {"default": "Default", "angel": "Angel", "devil": "Devil", "executor": "Executor"}
            role_name = role_names.get(role, role)

            route_msg = Message(
                platform="telegram",
                chat_id=chat_id,
                user_id=str(cb["from"]["id"]),
                text=f"/set {config_key} {model_name}",
            )
            result = self.route(route_msg)
            # Re-acknowledge with result text for feedback
            self._client.post(
                f"{self._api_base}/answerCallbackQuery",
                json={"callback_query_id": cb_id, "text": f"{role_name} → {model_name}"},
                timeout=5,
            )
            # Update button to show selection
            if msg_id:
                providers = self._get_providers_config()
                sel_provider = None
                for pn, mods in providers.items():
                    if model_name in mods:
                        sel_provider = pn
                        break
                if sel_provider:
                    self._send_model_buttons(chat_id, msg_id, sel_provider, model_name)
            if result:
                self.send(chat_id, result)

        elif data.startswith("court:"):
            # M5-D8: STAY-verdict inline keyboard callback. Re-derive
            # case_id + action from "court:{case_id}:{action}". Route
            # to /court resume which the dispatcher understands.
            try:
                _, case_id, action = data.split(":", 2)
            except ValueError:
                logger.warning(f"[Telegram] malformed court callback: {data!r}")
                return
            # Persist the user's decision on the case file (annotate the
            # archive JSON) so the agent can read it on resume. The
            # agent itself decides what to do with the action; here we
            # just acknowledge to Telegram and edit the original STAY
            # message to show "你揀咗: {action}".
            try:
                from pathlib import Path
                import json
                import time
                _case_path = Path.home() / ".baw" / "court" / "cases" / f"{case_id}.json"
                if _case_path.exists():
                    _case = json.loads(_case_path.read_text(encoding="utf-8"))
                    _case.setdefault("stay_decisions", []).append({
                        "action": action,
                        "user_id": str(cb.get("from", {}).get("id", "")),
                        "at": time.time(),
                    })
                    _case_path.write_text(
                        json.dumps(_case, ensure_ascii=False, indent=2),
                        encoding="utf-8",
                    )
            except Exception as _ce:
                logger.warning(f"[Telegram] could not persist stay decision: {_ce}")

            _labels = {
                "approve": "✅ 已批准執行",
                "backup":  "💾 將先備份再執行",
                "dismiss": "🚫 已撤案",
            }
            _ack_text = _labels.get(action, f"已選擇: {action}")
            self._client.post(
                f"{self._api_base}/answerCallbackQuery",
                json={"callback_query_id": cb_id, "text": _ack_text},
                timeout=5,
            )
            # Edit the original message to reflect the decision.
            if msg_id:
                try:
                    self._client.post(
                        f"{self._api_base}/editMessageText",
                        json={
                            "chat_id": chat_id,
                            "message_id": msg_id,
                            "text": f"⏸️ #{case_id} │ {_ack_text}\n"
                                     f"用戶 ID: <code>{cb.get('from', {}).get('id', '?')}</code>\n"
                                     f"決策已記錄。下一個 agent 回合會讀取呢個決定。",
                            "parse_mode": "HTML",
                        },
                        timeout=5,
                    )
                except Exception as _ee:
                    logger.warning(f"[Telegram] edit after STAY error: {_ee}")

        else:
            self._client.post(
                f"{self._api_base}/answerCallbackQuery",
                json={"callback_query_id": cb_id, "text": "Unknown"},
                timeout=5,
            )
